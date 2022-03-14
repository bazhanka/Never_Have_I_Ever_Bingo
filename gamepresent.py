from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import random


class GamePresent:
    def __init__(self, inch, input_file, output_file):
        self.inch = inch
        self.input_file = input_file
        self.output_file = output_file

    def generate_card(self):
        f = open(self.input_file, encoding='utf-8')
        fants = [line.strip() for line in f]
        all_keys = []
        a = 0
        new_phrases = []
        while a < self.inch:
            key = random.sample(range(0, len(fants)), 9)
            if key not in all_keys:
                all_keys.append(key)
                a += 1
                ph = []
                for el in key:
                    phrase = fants[el]
                    ph.append(phrase)
                new_phrases.append(ph)
                if a == self.inch:
                    return new_phrases

    def generate_table(self):
        prs = Presentation()
        prs.slide_height = Inches(3.7)
        prs.slide_width = Inches(5.43)
        count = 0
        all_keys = self.generate_card()
        all_index = 0
        while count < self.inch:
            title_only_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(title_only_slide_layout)
            shapes = slide.shapes
            rows = cols = 3
            left = top = Inches(0.0)
            width = Inches(5.43)
            height = Inches(3.7)
            table = shapes.add_table(rows, cols, left, top, width, height).table
            r = 0
            c = 0
            all_count = 0
            while c < 3 and r < 3:
                txt = table.cell(r, c).text_frame.paragraphs[0]
                txt.alignment = PP_ALIGN.CENTER
                table.cell(r, c).margin_left = Inches(0.2)
                table.cell(r, c).margin_top = Inches(0.2)
                fill = table.cell(r, c).fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 250, 250)
                run = txt.add_run()
                run.text = all_keys[all_index][all_count]
                font = run.font
                font.name = 'Verdana'
                font.size = Pt(10)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
                c += 1
                all_count += 1
                if c == 3 and all_count < 9:
                    r += 1
                    c = 0
                elif all_count == 9:
                    all_index += 1
                    count += 1
                    if count == self.inch:
                        prs.save(self.output_file)
                        break
